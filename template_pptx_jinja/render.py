import copy
import re

from jinja2 import exceptions, Environment
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.table import _Cell

from template_pptx_jinja import pictures
# 


class PPTXRendering:
    # Patrón regex para detectar placeholders del tipo {{ table:mi_tabla }}
    table_data_pattern = re.compile(r"\{\{\s*table\s*:\s*([^}\s]+)\s*\}\}")

    def __init__(self, input_path, data, output_path, env=None):
        # Ruta al archivo .pptx de plantilla
        self.input_path = input_path

        # Diccionario de datos que se inyectará en la plantilla
        self.model = data["model"]

        # Diccionario que asocia hashes de imagen con rutas a nuevas imágenes
        self.pictures = data.get("pictures")

        # Ruta de salida donde se guardará la presentación generada
        self.output_path = output_path

        # Entorno Jinja2 (puede pasarse uno externo o crear uno nuevo)
        self.env = env if env is not None else Environment()

        # Slide y shape actuales (para contexto interno)
        self.current_slide = None
        self.current_shape = None

        # Lista de mensajes de error encontrados durante el renderizado
        self.message_raw = []

    def process(self):
        # Carga la presentación base desde input_path
        ppt = Presentation(self.input_path)
        for slide in ppt.slides:
            self.current_slide = slide
            self._render_slide(slide)
        ppt.save(self.output_path)
        return "\n".join(self.message_raw)

    def _render_slide(self, slide):
        # Recorre cada forma (shape) en la diapositiva
        for shape in slide.shapes:
            self.current_shape = shape
            self._render_shape(shape)

    def _render_shape(self, shape):
        # Aplica el renderizado correspondiente según el tipo de contenido
        if shape.has_text_frame:
            self._render_text_frame(shape.text_frame)
        if shape.has_table:
            self._render_table(shape.table)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            self._render_picture(shape)

    def _render_picture(self, shape):
        # Compara el hash de la imagen actual con el de las imágenes a reemplazar
        for picture in self.pictures:
            if pictures.get_hash(picture) == shape.image.sha1:
                pictures.replace_img_slide(
                    self.current_slide, shape, self.pictures[picture]
                )

    def _render_text_frame(self, text_frame):
        for paragraph in text_frame.paragraphs:
            self._render_paragraph(paragraph)

    def _merge_placeholder_runs(self, paragraph):
        """
        Fusiona fragmentos de texto (runs) que forman una expresión {{ ... }}
        en uno solo para que Jinja2 pueda interpretarlo correctamente.
        """
        runs = paragraph.runs
        i = 0
        while i < len(runs):
            text = runs[i].text or ""
            if "{{" in text:
                merged = text
                j = i + 1
                while j < len(runs) and "}}" not in merged:
                    merged += runs[j].text or ""
                    j += 1
                if "}}" in merged:
                    runs[i].text = merged
                    for k in range(i + 1, j):
                        runs[k].text = ""
                    i = j
                else:
                    break
            else:
                i += 1

    def _render_paragraph(self, paragraph):
        # 1) Une los placeholders divididos entre runs
        self._merge_placeholder_runs(paragraph)
        # 2) Reemplaza los placeholders con Jinja2
        for run in paragraph.runs:
            self._render_run(run)

    def _render_table(self, table):
        # Prepara la tabla si tiene placeholders especiales (tipo {{ table:mi_tabla }})
        self._prepare_table(table)
        # Renderiza cada celda de la tabla
        for cell in table.iter_cells():
            self._render_cell(cell)

    def _prepare_table(self, table):
        for cell in table.iter_cells():
            matches = re.findall(self.table_data_pattern, cell.text)
            if matches:
                table_data_key = matches[0].strip()
                self._remove_row(table, 0)  # Elimina la fila de plantilla
                if table_data := self.model.get(table_data_key):
                    for i, row in enumerate(table_data):
                        self._add_row(table, i, table_data_key)
                    self._remove_row(table, 1)  # Elimina la fila vacía que queda
            break

    def _render_cell(self, cell):
        self._render_text_frame(cell.text_frame)

    def _render_run(self, run):
        try:
            template = self.env.from_string(str(run.text))
            rendered = template.render(self.model)
        except exceptions.UndefinedError as error:
            error_text = f"{error.__class__.__name__}: {error}"
            self.message_raw.append(error_text)
        except exceptions.TemplateSyntaxError as error:
            error_text = (
                f"{error.__class__.__name__}: {error}\n"
                "you should re-write the whole {{}} tag"
            )
            self.message_raw.append(error_text)
        else:
            run.text = rendered

    @staticmethod
    def _add_row(table, row_index, table_data_key):
        # Duplica la última fila como plantilla para insertar nuevos datos
        new_row = copy.deepcopy(table._tbl.tr_lst[-1])

        for cell_index, tc in enumerate(new_row.tc_lst):
            cell = _Cell(tc, new_row.tc_lst)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = f"{{{{ {table_data_key}[{row_index}][{cell_index}] }}}}"

        table._tbl.append(new_row)

    @staticmethod
    def _remove_row(table, row_index):
        # Elimina una fila del objeto tabla de PowerPoint
        row = list(table.rows)[row_index]
        table._tbl.remove(row._tr)
    # 
# 
